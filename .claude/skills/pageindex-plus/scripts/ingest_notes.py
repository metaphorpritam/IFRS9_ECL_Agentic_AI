# /// script
# requires-python = ">=3.10"
# dependencies = ["pypdf", "pymupdf", "beautifulsoup4", "lxml", "python-pptx",
#                 "openpyxl", "xlrd", "python-docx", "Pillow"]
# ///
"""Convert a folder of NOTES into a corpus for build_pageindex.py — TEXT *and* VISUALS.

For every supported document this writes one Markdown file per document AND extracts
embedded figures (raster images, vector charts, flowcharts, plots) as PNG assets. Each
figure is recorded in the Markdown as a **figure card** — a heading the index can anchor —
carrying the source file, the page / slide / sheet number, nearby-text keywords, any
in-image OCR text, and a relative path to the saved PNG. So a later tree-search finds the
figure by its surrounding context, and the answer can cite both the page and the asset.

    uv run ingest_notes.py <notes_dir> <corpus_dir> [options]

Supported: .pdf .html .htm .pptx .xlsx .xls (via xlrd, text/tables only) .docx
.md .markdown .txt
Visuals are extracted from: PDF (embedded rasters + rasterised vector-figure pages),
PPTX (slide pictures, tables, grouped textboxes, OMML equations — plus whole-slide
LibreOffice rasterisation for EMF/WMF/SmartArt/equation slides that python-pptx
cannot read), XLSX (embedded images + a flag for native charts), DOCX (inline
images). Markdown/HTML local <img>/![]() references are catalogued and copied
(refused if they resolve outside the sources tree).

INCREMENTAL: a manifest (`<corpus_dir>/.ingest_manifest.json`) records each source's
content hash. Re-running skips unchanged files, re-processes changed ones, and (with
--prune) drops outputs for deleted sources. Add sources over time and just re-run.

Key options:
  --assets-subdir img        where PNGs go, relative to corpus (default: img)
  --rasterize-dpi 150        DPI for PDF vector-figure page rasterisation
  --no-rasterize             skip whole-page rasterisation (embedded rasters only)
  --max-figures-per-doc 40   safety cap
  --ocr                      OCR extracted figures (needs pytesseract + tesseract)
  --prune                    delete corpus/asset outputs for sources no longer present
  --force                    ignore the manifest; reprocess everything

Captions are intentionally OFFLINE here (keywords + OCR). For richer descriptions, run
caption_figures.py afterwards (LLM-vision hook) — it merges captions in place without
re-extracting. Each handler imports its deps lazily, so one missing optional package
only disables that one file type or feature.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from pathlib import Path

SUPPORTED = {".pdf", ".html", ".htm", ".pptx", ".xlsx", ".xls", ".docx",
             ".md", ".markdown", ".txt"}
MANIFEST_NAME = ".ingest_manifest.json"
MIN_FIG_BYTES = 3000   # below this a PNG is almost certainly a mask/decoration
MIN_FIG_PX = 64        # drop tiny icons/bullets


# --------------------------------------------------------------------------- utils
def _slug(rel: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "_", rel).strip("_") or "doc"


def _hash_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()[:16]


def _keywords(text: str, k: int = 12) -> str:
    """A few distinctive nearby-text words to make the figure findable in tree-search."""
    words = re.findall(r"[A-Za-z][A-Za-z0-9_-]{2,}", text or "")
    seen, out = set(), []
    stop = {"the", "and", "for", "with", "this", "that", "from", "are", "was",
            "page", "slide", "sheet", "figure", "table", "see", "fig"}
    for w in words:
        lw = w.lower()
        if lw in stop or lw in seen:
            continue
        seen.add(lw)
        out.append(w)
        if len(out) >= k:
            break
    return " ".join(out)


def _ocr(png_path: Path) -> str:
    """Optional in-image text. Returns '' if OCR stack is unavailable."""
    try:
        import pytesseract  # noqa
        from PIL import Image
    except Exception:
        return ""
    try:
        txt = pytesseract.image_to_string(Image.open(png_path))
        return re.sub(r"\s+", " ", txt).strip()[:400]
    except Exception:
        return ""


def _save_png(data: bytes, dest: Path) -> bool:
    """Write PNG bytes (re-encoding via Pillow if needed). Drop tiny/empty images."""
    from PIL import Image
    import io
    try:
        im = Image.open(io.BytesIO(data))
        if im.width < MIN_FIG_PX and im.height < MIN_FIG_PX:
            return False
        if im.mode in ("CMYK", "P", "RGBA"):
            im = im.convert("RGB")
        dest.parent.mkdir(parents=True, exist_ok=True)
        im.save(dest, "PNG")
        return dest.stat().st_size >= MIN_FIG_BYTES or (im.width * im.height) > 50000
    except Exception as e:
        # EMF/WMF metafiles land here — Pillow cannot decode them. Never fail
        # silently (learnings §1/§7): the slide-rasterisation path is the recovery.
        if len(data) >= MIN_FIG_BYTES:
            print(f"  fig-skip ({dest.name}): undecodable image data "
                  f"({type(e).__name__}: {e}) — if this is a PPTX EMF/WMF metafile, "
                  f"slide rasterisation recovers it (--raster-slides)",
                  file=sys.stderr)
        return False


class FigCardWriter:
    """Accumulates figure cards (Markdown) and tracks asset files for the manifest."""
    def __init__(self, doc_slug: str, assets_dir: Path, assets_rel: str, do_ocr: bool):
        self.doc_slug = doc_slug
        self.assets_dir = assets_dir
        self.assets_rel = assets_rel
        self.do_ocr = do_ocr
        self.n = 0
        self.cards: list[str] = []
        self.assets: list[str] = []

    def add(self, png_bytes: bytes, *, anchor: str, kind: str, nearby: str) -> bool:
        self.n += 1
        name = f"{self.doc_slug}_fig{self.n:03d}.png"
        dest = self.assets_dir / name
        if not _save_png(png_bytes, dest):
            self.n -= 1
            return False
        rel = f"{self.assets_rel}/{name}"
        self.assets.append(str(dest))
        ocr = self._ocr_line(dest)
        kw = _keywords(nearby)
        card = [
            f"### Figure {self.n} ({anchor})",
            "",
            f"- asset: `{rel}`",
            f"- kind: {kind}",
            f"- anchor: {anchor}",
        ]
        if kw:
            card.append(f"- context_keywords: {kw}")
        if ocr:
            card.append(f"- ocr_text: {ocr}")
        card.append("- caption: (uncaptioned — run caption_figures.py)")
        self.cards.append("\n".join(card))
        return True

    def _ocr_line(self, dest: Path) -> str:
        return self._ocr_cache.setdefault(dest, _ocr(dest)) if self.do_ocr else ""

    _ocr_cache: dict = {}

    def markdown(self) -> str:
        if not self.cards:
            return ""
        return "\n\n## Figures\n\n" + "\n\n".join(self.cards)


# --------------------------------------------------------------------------- handlers
# Each returns (text_markdown, fig_writer). fig_writer may have 0 figures.

def from_pdf(path, slug, assets_dir, assets_rel, args):
    from pypdf import PdfReader
    fw = FigCardWriter(slug, assets_dir, assets_rel, args.ocr)
    out = []
    reader = PdfReader(str(path))
    page_text = []
    for i, page in enumerate(reader.pages, 1):
        txt = (page.extract_text() or "").strip()
        page_text.append(txt)
        if txt:
            out.append(f"## Page {i}\n\n{txt}")

    # Visuals via PyMuPDF: embedded rasters always; whole-page rasterise pages that
    # likely hold a vector chart/flowchart (little text, or drawing ops present).
    try:
        import fitz
    except Exception:
        print("  (pymupdf missing — PDF text only)", file=sys.stderr)
        return "\n\n".join(out), fw

    doc = fitz.open(str(path))
    seen_xref = set()
    for i, pg in enumerate(doc, 1):
        if fw.n >= args.max_figures_per_doc:
            break
        nearby = page_text[i - 1] if i - 1 < len(page_text) else ""
        # 1) embedded raster images
        for img in pg.get_images(full=True):
            xref = img[0]
            if xref in seen_xref:
                continue
            seen_xref.add(xref)
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha >= 4:        # CMYK etc -> RGB
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                if fw.add(pix.tobytes("png"), anchor=f"{path.name} p{i}",
                          kind="pdf-raster", nearby=nearby):
                    pass
            except Exception:
                continue
            if fw.n >= args.max_figures_per_doc:
                break
        # 2) vector figure page -> rasterise the page (catches matplotlib/Excel/draw.io)
        if not args.no_rasterize:
            has_drawings = bool(pg.get_drawings())
            sparse_text = len((page_text[i - 1] or "")) < 200
            if has_drawings and sparse_text:
                try:
                    zoom = args.rasterize_dpi / 72.0
                    pix = pg.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                    fw.add(pix.tobytes("png"), anchor=f"{path.name} p{i} (page render)",
                           kind="pdf-vector-page", nearby=nearby)
                except Exception:
                    pass
    doc.close()
    return "\n\n".join(out), fw


def from_html(path, slug, assets_dir, assets_rel, args):
    from bs4 import BeautifulSoup
    fw = FigCardWriter(slug, assets_dir, assets_rel, args.ocr)
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "lxml")
    # catalogue local images (copy referenced files that exist on disk)
    for img in soup.find_all("img"):
        src = (img.get("src") or "").strip()
        if not src or src.startswith(("http://", "https://", "data:")):
            continue
        ref = (path.parent / src).resolve()
        root = getattr(args, "notes_root", None)
        if root and not ref.is_relative_to(root):
            print(f"  refuse ({path.name}): <img> ref escapes the sources tree, "
                  f"not copied: {src}", file=sys.stderr)
            continue
        if ref.exists() and ref.suffix.lower() in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
            try:
                fw.add(ref.read_bytes(), anchor=f"{path.name} <img {src}>",
                       kind="html-img", nearby=img.get("alt", "") + " " + src)
            except Exception as e:
                print(f"  warn ({path.name}): <img {src}> unreadable — "
                      f"{type(e).__name__}: {e}", file=sys.stderr)
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    for h in soup.find_all(re.compile(r"^h[1-6]$")):
        h.insert_before("\n" + "#" * int(h.name[1]) + " ")
        h.insert_after("\n")
    text = re.sub(r"\n{3,}", "\n\n", soup.get_text("\n")).strip()
    return text, fw


# ---- PPTX deep extraction (learnings §0/§1) ---------------------------------
# python-pptx reads text frames and raster pictures ONLY. Four content classes
# are invisible to a flat shapes loop and were silently dropped before:
#   tables, grouped shapes, OMML (Office Math) equations, EMF/WMF metafiles.
# The first three are recovered below by walking the XML; the fourth (plus
# SmartArt and anything vector) is recovered by rasterising the slide via
# LibreOffice headless -> PDF -> PyMuPDF PNG.

_VECTOR_PIC_EXTS = {"emf", "wmf", "emz", "wmz", "pict"}
# OMML spans, prefix-agnostic: real decks use m:, but lxml re-serialisation may
# assign any prefix (or none) — match whatever prefix opens the oMath element.
_OMML_SPAN = re.compile(r"<(?:(\w+):)?oMath(?:\s[^>]*)?>(.*?)</(?:\1:)?oMath>",
                        re.S)


def _omml_text(xml: str) -> str:
    """Concatenate the <t> runs inside every oMath span, using each span's own
    namespace prefix."""
    parts = []
    for pref, body in _OMML_SPAN.findall(xml):
        p = re.escape(pref + ":") if pref else ""
        parts += re.findall(r"<" + p + r"t(?:\s[^>]*)?>([^<]*)</" + p + r"t>",
                            body)
    return "".join(parts).strip()


def _soffice_exe():
    """Locate LibreOffice. Env var SOFFICE wins; then PATH; then default installs."""
    import os
    import shutil
    cands = [os.environ.get("SOFFICE"), "soffice",
             r"C:\Program Files\LibreOffice\program\soffice.exe",
             r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
             "/usr/bin/soffice",
             "/Applications/LibreOffice.app/Contents/MacOS/soffice"]
    for c in cands:
        if not c:
            continue
        hit = shutil.which(c) or (c if Path(c).exists() else None)
        if hit:
            return hit
    return None


def _rasterize_pptx_slides(path, slide_nos, dpi, warns):
    """PPTX -> PDF (soffice --headless) -> {slide_no: png_bytes} via PyMuPDF.

    Renders EMF/WMF/OMML/SmartArt/tables exactly as they appear on the slide.
    A MuPDF `No common ancestor in structure tree` warning here is harmless.
    """
    out = {}
    if not slide_nos:
        return out
    exe = _soffice_exe()
    if not exe:
        warns.append("LibreOffice (soffice) NOT FOUND — slides "
                     f"{sorted(slide_nos)} hold equations/SmartArt/metafiles that "
                     "python-pptx cannot read and were NOT rasterised. Install "
                     "LibreOffice or set SOFFICE=<path to soffice(.exe)>.")
        return out
    try:
        import fitz
    except Exception as e:
        warns.append(f"pymupdf unavailable ({e}) — cannot rasterise slides")
        return out
    import subprocess
    import tempfile
    with tempfile.TemporaryDirectory() as td:
        try:
            subprocess.run([exe, "--headless", "--convert-to", "pdf",
                            "--outdir", td, str(path)],
                           check=True, capture_output=True, timeout=300)
        except Exception as e:
            warns.append(f"soffice conversion failed — {type(e).__name__}: {e}")
            return out
        pdfs = sorted(Path(td).glob("*.pdf"))
        if not pdfs:
            warns.append("soffice reported success but produced no PDF")
            return out
        doc = fitz.open(str(pdfs[0]))
        zoom = dpi / 72.0
        for n in sorted(slide_nos):
            if 1 <= n <= doc.page_count:
                try:
                    pix = doc[n - 1].get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                    out[n] = pix.tobytes("png")
                except Exception as e:
                    warns.append(f"slide {n}: rasterise failed — "
                                 f"{type(e).__name__}: {e}")
            else:
                warns.append(f"slide {n}: beyond converted PDF "
                             f"({doc.page_count} pages)")
        doc.close()
    return out


def _walk_pptx_shapes(shapes, slide_no, warns):
    """Depth-first shape walk, recursing into GROUPs (learnings §1).

    Yields (kind, payload): kind in {'text','table','equation','picture',
    'chart','needs_raster'}. Every extraction failure lands in `warns` —
    nothing is swallowed silently.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        # -- groups: recurse ---------------------------------------------------
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _walk_pptx_shapes(sh.shapes, slide_no, warns)
                continue
        except Exception as e:
            warns.append(f"slide {slide_no}: group probe — {type(e).__name__}: {e}")
        # -- tables (CAPM series, case exhibits, ...) ---------------------------
        try:
            if getattr(sh, "has_table", False):
                t = sh.table
                rows = [" | ".join((c.text or "").strip().replace("\n", " ")
                                   for c in r.cells) for r in t.rows]
                yield "table", rows
                yield "needs_raster", "table"     # render layout too
                continue
        except Exception as e:
            warns.append(f"slide {slide_no}: table read — {type(e).__name__}: {e}")
        # -- OMML equations live in raw shape XML, never in text frames ---------
        got_equation = False
        try:
            xml = sh._element.xml
            if "oMath" in xml:
                runs = _omml_text(xml)
                yield "equation", runs or "(equation present — see slide render)"
                yield "needs_raster", "OMML equation"
                got_equation = True
            if "drawingml/2006/diagram" in xml:     # SmartArt
                yield "needs_raster", "SmartArt"
        except Exception as e:
            warns.append(f"slide {slide_no}: xml sniff — {type(e).__name__}: {e}")
        # -- plain text frames ---------------------------------------------------
        try:
            if sh.has_text_frame:
                txt = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
                if txt and not got_equation:
                    yield "text", (sh, txt)
        except Exception as e:
            warns.append(f"slide {slide_no}: text read — {type(e).__name__}: {e}")
        # -- pictures (EMF/WMF flagged for rasterisation) ------------------------
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ext = (sh.image.ext or "").lower()
                yield "picture", sh
                if ext in _VECTOR_PIC_EXTS:
                    yield "needs_raster", f"{ext} metafile"
        except Exception as e:
            warns.append(f"slide {slide_no}: picture read — {type(e).__name__}: {e}")
        try:
            if getattr(sh, "has_chart", False):
                yield "chart", sh
        except Exception as e:
            warns.append(f"slide {slide_no}: chart probe — {type(e).__name__}: {e}")


def from_pptx(path, slug, assets_dir, assets_rel, args):
    from pptx import Presentation
    fw = FigCardWriter(slug, assets_dir, assets_rel, args.ocr)
    prs = Presentation(str(path))
    out = []
    warns: list[str] = []
    raster_wanted: dict[int, list[str]] = {}
    slide_kw: dict[int, str] = {}
    mode = getattr(args, "raster_slides", "auto")

    i = 0
    for i, slide in enumerate(prs.slides, 1):
        title, body, kw_pool = "", [], []
        for kind, payload in _walk_pptx_shapes(slide.shapes, i, warns):
            if kind == "text":
                sh, txt = payload
                kw_pool.append(txt)
                try:
                    is_title = (sh == slide.shapes.title)
                except Exception:
                    is_title = False
                if not title and is_title:
                    title = txt
                else:
                    body.append(txt)
            elif kind == "table":
                kw_pool.extend(payload[:2])
                body.append("[TABLE]\n" + "\n".join("| " + r + " |"
                                                     for r in payload))
            elif kind == "equation":
                kw_pool.append(payload)
                body.append(f"[EQUATION] {payload}")
            elif kind == "picture":
                try:
                    fw.add(payload.image.blob, anchor=f"{path.name} slide {i}",
                           kind="pptx-picture",
                           nearby=" ".join(kw_pool) + " " + title)
                except Exception as e:
                    warns.append(f"slide {i}: picture save — "
                                 f"{type(e).__name__}: {e}")
            elif kind == "chart":
                fw.cards.append(
                    f"### Figure {fw.n + 1} ({path.name} slide {i})\n\n"
                    f"- kind: pptx-native-chart\n- anchor: {path.name} slide {i}\n"
                    f"- context_keywords: {_keywords(' '.join(kw_pool))}\n"
                    f"- note: native chart (no raster); regenerate from data at "
                    f"notes time\n"
                    f"- caption: (uncaptioned — run caption_figures.py)")
                fw.n += 1
                raster_wanted.setdefault(i, []).append("native chart")
            elif kind == "needs_raster":
                raster_wanted.setdefault(i, []).append(payload)

        head = f"## Slide {i}" + (f": {title}" if title else "")
        out.append(head + ("\n\n" + "\n\n".join(body) if body else ""))
        slide_kw[i] = (title + " " + " ".join(kw_pool))[:300]
        if fw.n >= args.max_figures_per_doc:
            warns.append(f"figure cap ({args.max_figures_per_doc}) reached at "
                         f"slide {i}; later slides not scanned for figures")
            break

    # -- whole-slide rasterisation (learnings §0) -----------------------------
    if mode == "all":
        targets = set(range(1, i + 1))
    elif mode == "off":
        targets = set()
    else:  # auto
        targets = set(raster_wanted)
    budget = max(0, args.max_figures_per_doc - fw.n)
    for n, png in _rasterize_pptx_slides(path, set(sorted(targets)[:budget]),
                                         args.rasterize_dpi, warns).items():
        why = ", ".join(sorted(set(raster_wanted.get(n, ["requested"]))))
        fw.add(png, anchor=f"{path.name} slide {n} (page render)",
               kind="pptx-slide-render",
               nearby=(slide_kw.get(n, "") + " " + why).strip())

    for w in warns:
        print(f"  warn ({path.name}): {w}", file=sys.stderr)
    return "\n\n".join(out), fw


def from_excel(path, slug, assets_dir, assets_rel, args):
    if path.suffix.lower() == ".xls":          # BIFF — openpyxl raises on these
        return _from_xls_legacy(path, slug, assets_dir, assets_rel, args)
    from openpyxl import load_workbook
    fw = FigCardWriter(slug, assets_dir, assets_rel, args.ocr)
    wb = load_workbook(str(path), read_only=False, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"## Sheet: {ws.title}")
        rows = []
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r >= 60:
                rows.append("| ... |")
                break
            cells = [("" if v is None else str(v)).replace("|", "\\|")[:40] for v in row[:12]]
            rows.append("| " + " | ".join(cells) + " |")
            if r == 0:
                rows.append("| " + " | ".join("---" for _ in cells) + " |")
        out.append("\n".join(rows))
        # embedded images on the sheet
        for im in getattr(ws, "_images", []):
            try:
                blob = im._data() if callable(getattr(im, "_data", None)) else im.ref
                if isinstance(blob, (bytes, bytearray)):
                    fw.add(bytes(blob), anchor=f"{path.name} sheet {ws.title}",
                           kind="xlsx-image", nearby=ws.title)
            except Exception:
                pass
        # native charts -> record a card; data is recoverable from the sheet table
        for _ch in getattr(ws, "_charts", []):
            fw.cards.append(
                f"### Figure {fw.n + 1} ({path.name} sheet {ws.title})\n\n"
                f"- kind: xlsx-native-chart\n- anchor: {path.name} sheet {ws.title}\n"
                f"- context_keywords: {_keywords(ws.title)}\n"
                f"- note: native chart (no raster); data is in the sheet table above — "
                f"regenerate cleanly in matplotlib at notes time\n"
                f"- caption: (uncaptioned — run caption_figures.py)")
            fw.n += 1
    return "\n\n".join(out), fw


def _from_xls_legacy(path, slug, assets_dir, assets_rel, args):
    """Legacy BIFF .xls via xlrd (learnings §4). openpyxl only reads OOXML .xlsx —
    before this fallback every .xls was declared supported but silently skipped.
    Text/tables only: xlrd exposes no embedded-image API, and that limitation is
    REPORTED below rather than hidden."""
    import xlrd
    fw = FigCardWriter(slug, assets_dir, assets_rel, args.ocr)
    wb = xlrd.open_workbook(str(path))
    out = []
    for ws in wb.sheets():
        out.append(f"## Sheet: {ws.name}")
        rows = []
        for r in range(min(ws.nrows, 60)):
            cells = []
            for c in range(min(ws.ncols, 12)):
                cell = ws.cell(r, c)
                v = cell.value
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        v = xlrd.xldate.xldate_as_datetime(
                            v, wb.datemode).isoformat(sep=" ")
                    except Exception:
                        pass
                cells.append(("" if v in ("", None) else str(v))
                             .replace("|", "\\|")[:40])
            rows.append("| " + " | ".join(cells) + " |")
            if r == 0:
                rows.append("| " + " | ".join("---" for _ in cells) + " |")
        if ws.nrows > 60:
            rows.append("| ... |")
        out.append("\n".join(rows))
    print(f"  note ({path.name}): legacy .xls read via xlrd — text/tables only; "
          f"BIFF embedded images are not extractable (re-save as .xlsx to "
          f"recover them)", file=sys.stderr)
    return "\n\n".join(out), fw


def from_docx(path, slug, assets_dir, assets_rel, args):
    from docx import Document
    fw = FigCardWriter(slug, assets_dir, assets_rel, args.ocr)
    doc = Document(str(path))
    out = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        style = (p.style.name or "").lower() if p.style else ""
        m = re.search(r"heading (\d)", style)
        out.append(("#" * min(int(m.group(1)) + 1, 6) + " " + t) if m else t)
    body = "\n".join(out)
    # inline images live in the docx zip part store
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                blob = rel.target_part.blob
                fw.add(blob, anchor=f"{path.name} inline image",
                       kind="docx-image", nearby=body[:300])
                if fw.n >= args.max_figures_per_doc:
                    break
    except Exception:
        pass
    return body, fw


def from_plain(path, slug, assets_dir, assets_rel, args):
    """Markdown / txt: keep text; catalogue & copy local image references in Markdown."""
    fw = FigCardWriter(slug, assets_dir, assets_rel, args.ocr)
    text = path.read_text(encoding="utf-8", errors="replace")
    if path.suffix.lower() in {".md", ".markdown"}:
        for m in re.finditer(r"!\[[^\]]*\]\(([^)]+)\)", text):
            src = m.group(1).split()[0].strip("<>")
            if src.startswith(("http://", "https://", "data:")):
                continue
            ref = (path.parent / src).resolve()
            root = getattr(args, "notes_root", None)
            if root and not ref.is_relative_to(root):
                print(f"  refuse ({path.name}): ![]() ref escapes the sources "
                      f"tree, not copied: {src}", file=sys.stderr)
                continue
            if ref.exists() and ref.suffix.lower() in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
                try:
                    fw.add(ref.read_bytes(), anchor=f"{path.name} ![]({src})",
                           kind="md-img", nearby=text[max(0, m.start() - 120):m.start()])
                except Exception as e:
                    print(f"  warn ({path.name}): ![]({src}) unreadable — "
                          f"{type(e).__name__}: {e}", file=sys.stderr)
    return text, fw


HANDLERS = {
    ".pdf": from_pdf, ".html": from_html, ".htm": from_html, ".pptx": from_pptx,
    ".xlsx": from_excel, ".xls": from_excel, ".docx": from_docx,
    ".md": from_plain, ".markdown": from_plain, ".txt": from_plain,
}


# --------------------------------------------------------------------------- manifest
def _load_manifest(corpus: Path) -> dict:
    mf = corpus / MANIFEST_NAME
    if mf.exists():
        try:
            return json.loads(mf.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"files": {}}


def _save_manifest(corpus: Path, manifest: dict) -> None:
    (corpus / MANIFEST_NAME).write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")


# --------------------------------------------------------------------------- main
def main(argv: list[str] | None = None) -> int:
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("notes_dir")
    ap.add_argument("corpus_dir")
    ap.add_argument("--assets-subdir", default="img")
    ap.add_argument("--rasterize-dpi", type=int, default=150)
    ap.add_argument("--no-rasterize", action="store_true")
    ap.add_argument("--max-figures-per-doc", type=int, default=40)
    ap.add_argument("--raster-slides", choices=["auto", "all", "off"],
                    default="auto",
                    help="LibreOffice whole-slide rasterisation for PPTX: auto = "
                         "only slides holding OMML equations / SmartArt / EMF-WMF "
                         "metafiles / tables / native charts; all; off")
    ap.add_argument("--ocr", action="store_true")
    ap.add_argument("--prune", action="store_true")
    ap.add_argument("--force", action="store_true")
    args = ap.parse_args(argv)

    notes = Path(args.notes_dir)
    args.notes_root = notes.resolve()   # containment root for referenced images
    corpus = Path(args.corpus_dir)
    assets_dir = corpus / args.assets_subdir
    corpus.mkdir(parents=True, exist_ok=True)
    assets_dir.mkdir(parents=True, exist_ok=True)

    manifest = {"files": {}} if args.force else _load_manifest(corpus)
    files = [p for p in sorted(notes.rglob("*"))
             if p.is_file() and p.suffix.lower() in SUPPORTED]

    ok = skip = unchanged = fig_total = 0
    current_rels = set()
    for p in files:
        rel = p.relative_to(notes).as_posix()
        current_rels.add(rel)
        slug = _slug(rel)
        digest = _hash_file(p)
        prev = manifest["files"].get(rel)
        if prev and prev.get("hash") == digest and not args.force:
            unchanged += 1
            fig_total += prev.get("figures", len(prev.get("assets", [])))
            continue
        # changed/new: clear old assets for this source first
        if prev:
            for a in prev.get("assets", []):
                try:
                    Path(a).unlink(missing_ok=True)
                except Exception:
                    pass
        try:
            handler = HANDLERS[p.suffix.lower()]
            text, fw = handler(p, slug, assets_dir, args.assets_subdir, args)
            text = (text or "").strip()
            fig_md = fw.markdown()
            if not text and not fig_md:
                print(f"  skip (empty): {rel}", file=sys.stderr)
                skip += 1
                manifest["files"].pop(rel, None)
                continue
            (corpus / f"{slug}.md").write_text(
                f"# {rel}\n\n{text}\n{fig_md}\n", encoding="utf-8")
            manifest["files"][rel] = {"hash": digest, "slug": slug,
                                      "assets": fw.assets, "figures": fw.n}
            ok += 1
            fig_total += fw.n
            print(f"  ok: {rel}  ({fw.n} figure(s))")
        except ImportError as e:
            print(f"  skip ({rel}): missing dependency — {e}", file=sys.stderr)
            skip += 1
        except Exception as e:  # noqa: BLE001
            print(f"  skip ({rel}): {type(e).__name__}: {e}", file=sys.stderr)
            skip += 1

    # prune deleted sources
    pruned = 0
    if args.prune:
        for rel in list(manifest["files"].keys()):
            if rel not in current_rels:
                rec = manifest["files"].pop(rel)
                Path(corpus / f"{rec.get('slug', _slug(rel))}.md").unlink(missing_ok=True)
                for a in rec.get("assets", []):
                    Path(a).unlink(missing_ok=True)
                pruned += 1

    _save_manifest(corpus, manifest)
    others = [p for p in sorted(notes.rglob("*"))
              if p.is_file() and p.suffix.lower() not in SUPPORTED]
    print(f"\ningested {ok} new/changed, {unchanged} unchanged -> {corpus}  "
          f"({skip} skipped, {pruned} pruned, {len(others)} unsupported, "
          f"{fig_total} figures total)")
    if ok == 0 and unchanged == 0:
        print("no documents ingested", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
